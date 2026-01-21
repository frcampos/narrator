#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
pptx_handler.py - Manipulacao de ficheiros PowerPoint
"""

import os
import shutil
import tempfile
from pathlib import Path
from typing import List, Optional, Tuple
from dataclasses import dataclass, field

# python-pptx
try:
    from pptx import Presentation
    from pptx.util import Cm, Emu
    PPTX_DISPONIVEL = True
except ImportError:
    PPTX_DISPONIVEL = False

from PIL import Image, ImageDraw


@dataclass
class SlideInfo:
    """InformaÃ§Ã£o de um slide"""
    numero: int
    texto_visivel: str = ""
    notas: str = ""
    texto_narrar: str = ""
    caminho_audio: str = ""
    duracao_audio: float = 0.0
    # Campos para traduÃ§Ã£o
    texto_traduzido: str = ""
    caminho_audio_traduzido: str = ""
    duracao_audio_traduzido: float = 0.0
    
    def tem_texto(self) -> bool:
        return bool(self.texto_narrar.strip())
    
    def tem_traducao(self) -> bool:
        return bool(self.texto_traduzido.strip())


@dataclass
class ConfigIcone:
    """ConfiguraÃ§Ã£o do Ã­cone de Ã¡udio"""
    mostrar: bool = True
    # PosiÃ§Ãµes: sup_dir, sup_esq, inf_dir, inf_esq (cantos)
    #           centro_sup, centro_inf, centro_esq, centro_dir (laterais centrais)
    posicao: str = "sup_dir"
    tamanho_cm: float = 1.0


@dataclass
class ApresentacaoInfo:
    """InformaÃ§Ã£o completa de uma apresentaÃ§Ã£o"""
    caminho: str
    nome: str = ""
    slides: List[SlideInfo] = field(default_factory=list)
    
    def __post_init__(self):
        if not self.nome:
            self.nome = Path(self.caminho).stem


class GestorPPTX:
    """Gestor de ficheiros PowerPoint"""
    
    def __init__(self):
        self.apresentacao: Optional[ApresentacaoInfo] = None
        self.config_icone = ConfigIcone()
    
    @staticmethod
    def disponivel() -> bool:
        """Verifica se python-pptx estÃ¡ disponÃ­vel"""
        return PPTX_DISPONIVEL
    
    # =========================================================================
    # EXTRAÃ‡ÃƒO DE TEXTO
    # =========================================================================
    
    def _extrair_texto_slide(self, slide) -> str:
        """Extrai texto visÃ­vel de um slide"""
        textos = []
        
        # Primeiro tentar tÃ­tulo e subtÃ­tulo (slides de tÃ­tulo)
        if slide.shapes.title:
            titulo = slide.shapes.title.text.strip()
            if titulo:
                textos.append(titulo)
        
        # Procurar placeholder de subtÃ­tulo
        for shape in slide.shapes:
            if shape.is_placeholder:
                if hasattr(shape, 'placeholder_format'):
                    ph_type = shape.placeholder_format.type
                    # Tipos 1=tÃ­tulo, 2=corpo, 3=subtÃ­tulo
                    if ph_type == 2 or ph_type == 3:  # corpo ou subtÃ­tulo
                        if shape.has_text_frame:
                            texto = shape.text_frame.text.strip()
                            if texto and texto not in textos:
                                textos.append(texto)
        
        # Depois outras shapes com texto
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    texto = para.text.strip()
                    if texto and texto not in textos:
                        textos.append(texto)
        
        return " ".join(textos)
    
    def _extrair_notas_slide(self, slide) -> str:
        """Extrai notas do orador"""
        try:
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                return notes_frame.text.strip() if notes_frame else ""
        except:
            pass
        return ""
    
    def abrir(self, caminho: str) -> bool:
        """
        Abre uma apresentaÃ§Ã£o PowerPoint.
        
        Returns:
            True se aberto com sucesso
        """
        if not PPTX_DISPONIVEL:
            return False
        
        if not os.path.exists(caminho):
            return False
        
        try:
            prs = Presentation(caminho)
            
            slides = []
            for i, slide in enumerate(prs.slides):
                texto_visivel = self._extrair_texto_slide(slide)
                notas = self._extrair_notas_slide(slide)
                
                # Por padrÃ£o, usar notas se existirem, senÃ£o texto do slide
                texto_narrar = notas if notas else texto_visivel
                
                slides.append(SlideInfo(
                    numero=i + 1,
                    texto_visivel=texto_visivel,
                    notas=notas,
                    texto_narrar=texto_narrar
                ))
            
            self.apresentacao = ApresentacaoInfo(
                caminho=caminho,
                slides=slides
            )
            
            return True
            
        except Exception as e:
            print(f"Erro ao abrir PPTX: {e}")
            return False
    
    @property
    def num_slides(self) -> int:
        """NÃºmero de slides"""
        if self.apresentacao:
            return len(self.apresentacao.slides)
        return 0
    
    def obter_slide(self, numero: int) -> Optional[SlideInfo]:
        """ObtÃ©m informaÃ§Ã£o de um slide (1-indexed)"""
        if self.apresentacao and 1 <= numero <= self.num_slides:
            return self.apresentacao.slides[numero - 1]
        return None
    
    def atualizar_texto_narrar(self, numero: int, texto: str):
        """Atualiza o texto a narrar de um slide"""
        slide = self.obter_slide(numero)
        if slide:
            slide.texto_narrar = texto
    
    def atualizar_texto_traduzido(self, numero: int, texto: str):
        """Atualiza o texto traduzido de um slide"""
        slide = self.obter_slide(numero)
        if slide:
            slide.texto_traduzido = texto
    
    def definir_audio_slide(self, numero: int, caminho: str, duracao: float):
        """Define o Ã¡udio gerado para um slide"""
        slide = self.obter_slide(numero)
        if slide:
            slide.caminho_audio = caminho
            slide.duracao_audio = duracao
    
    def definir_audio_traduzido_slide(self, numero: int, caminho: str, duracao: float):
        """Define o Ã¡udio traduzido para um slide"""
        slide = self.obter_slide(numero)
        if slide:
            slide.caminho_audio_traduzido = caminho
            slide.duracao_audio_traduzido = duracao
    
    # =========================================================================
    # CRIAÃ‡ÃƒO DE ÃCONE
    # =========================================================================
    
    def _criar_icone_audio(self, tamanho: int = 64, cor: tuple = (41, 128, 185, 255)) -> str:
        """Cria Ã­cone de som em PNG"""
        img = Image.new('RGBA', (tamanho, tamanho), (0, 0, 0, 0))
        draw = ImageDraw.Draw(img)
        
        # Altifalante
        pontos = [
            (tamanho * 0.12, tamanho * 0.35),
            (tamanho * 0.30, tamanho * 0.35),
            (tamanho * 0.50, tamanho * 0.12),
            (tamanho * 0.50, tamanho * 0.88),
            (tamanho * 0.30, tamanho * 0.65),
            (tamanho * 0.12, tamanho * 0.65),
        ]
        draw.polygon(pontos, fill=cor)
        
        # Ondas de som
        for r in [0.20, 0.32, 0.44]:
            bbox = [
                tamanho * 0.48,
                tamanho * (0.5 - r),
                tamanho * (0.48 + r * 1.2),
                tamanho * (0.5 + r)
            ]
            draw.arc(bbox, -55, 55, fill=cor, width=max(2, tamanho // 20))
        
        caminho = tempfile.mktemp(suffix='.png')
        img.save(caminho, 'PNG')
        return caminho
    
    def _criar_icone_traducao(self, tamanho: int = 64, cor: tuple = (39, 174, 96, 255)) -> str:
        """Cria Ã­cone de traduÃ§Ã£o (globo com texto) em PNG"""
        img = Image.new('RGBA', (tamanho, tamanho), (0, 0, 0, 0))
        draw = ImageDraw.Draw(img)
        
        # CÃ­rculo (globo)
        margem = tamanho * 0.1
        draw.ellipse([margem, margem, tamanho - margem, tamanho - margem], 
                     outline=cor, width=max(2, tamanho // 16))
        
        # Linhas do globo (meridianos)
        centro = tamanho / 2
        raio = (tamanho - 2 * margem) / 2
        # Elipse vertical
        draw.ellipse([centro - raio * 0.4, margem, centro + raio * 0.4, tamanho - margem],
                     outline=cor, width=max(1, tamanho // 24))
        # Linha horizontal
        draw.line([(margem, centro), (tamanho - margem, centro)], fill=cor, width=max(1, tamanho // 24))
        
        # Letra "A" pequeno
        try:
            from PIL import ImageFont
            fonte = ImageFont.truetype("arial.ttf", int(tamanho * 0.25))
        except:
            fonte = ImageFont.load_default()
        draw.text((tamanho * 0.65, tamanho * 0.55), "A", fill=cor, font=fonte)
        
        caminho = tempfile.mktemp(suffix='.png')
        img.save(caminho, 'PNG')
        return caminho
    
    def _calcular_posicao_icone(self, prs, incluir_segundo: bool = False, tem_legenda: bool = False) -> Tuple[Tuple[int, int], Tuple[int, int]]:
        """
        Calcula posiÃ§Ã£o dos Ã­cones baseado na configuraÃ§Ã£o.
        Quando hÃ¡ dois Ã­cones, ajusta para ambos ficarem visÃ­veis.
        
        Args:
            prs: ApresentaÃ§Ã£o
            incluir_segundo: Se hÃ¡ segundo Ã­cone (traduÃ§Ã£o)
            tem_legenda: Se hÃ¡ caixa de texto de legenda (ajusta posiÃ§Ã£o inferior)
        
        Returns:
            Tuplo com (pos_icone1, pos_icone2)
        """
        largura = prs.slide_width
        altura = prs.slide_height
        tamanho = Cm(self.config_icone.tamanho_cm)
        margem = Cm(0.5)
        espaco_entre = Cm(0.3)
        
        # Se tem legenda na parte inferior, subir os Ã­cones inferiores
        margem_inferior = Cm(2.5) if tem_legenda else margem
        
        # Largura total necessÃ¡ria para dois Ã­cones
        largura_dois = tamanho * 2 + espaco_entre
        
        # Centro do slide
        centro_x = (largura - tamanho) // 2
        centro_y = (altura - tamanho) // 2
        
        posicao = self.config_icone.posicao
        
        # Cantos
        if posicao == "sup_dir":
            if incluir_segundo:
                pos1 = (largura - largura_dois - margem, margem)
                pos2 = (largura - tamanho - margem, margem)
            else:
                pos1 = (largura - tamanho - margem, margem)
                pos2 = pos1
        elif posicao == "sup_esq":
            pos1 = (margem, margem)
            pos2 = (margem + tamanho + espaco_entre, margem)
        elif posicao == "inf_dir":
            y_pos = altura - tamanho - margem_inferior
            if incluir_segundo:
                pos1 = (largura - largura_dois - margem, y_pos)
                pos2 = (largura - tamanho - margem, y_pos)
            else:
                pos1 = (largura - tamanho - margem, y_pos)
                pos2 = pos1
        elif posicao == "inf_esq":
            y_pos = altura - tamanho - margem_inferior
            pos1 = (margem, y_pos)
            pos2 = (margem + tamanho + espaco_entre, y_pos)
        
        # PosiÃ§Ãµes centrais (laterais)
        elif posicao == "centro_esq":
            # Centro da margem esquerda
            pos1 = (margem, centro_y)
            pos2 = (margem, centro_y + tamanho + espaco_entre)
        elif posicao == "centro_dir":
            # Centro da margem direita
            if incluir_segundo:
                pos1 = (largura - tamanho - margem, centro_y - tamanho // 2 - espaco_entre // 2)
                pos2 = (largura - tamanho - margem, centro_y + tamanho // 2 + espaco_entre // 2)
            else:
                pos1 = (largura - tamanho - margem, centro_y)
                pos2 = pos1
        elif posicao == "centro_sup":
            # Centro superior
            if incluir_segundo:
                pos1 = (centro_x - tamanho // 2 - espaco_entre // 2, margem)
                pos2 = (centro_x + tamanho // 2 + espaco_entre // 2, margem)
            else:
                pos1 = (centro_x, margem)
                pos2 = pos1
        elif posicao == "centro_inf":
            # Centro inferior
            y_pos = altura - tamanho - margem_inferior
            if incluir_segundo:
                pos1 = (centro_x - tamanho // 2 - espaco_entre // 2, y_pos)
                pos2 = (centro_x + tamanho // 2 + espaco_entre // 2, y_pos)
            else:
                pos1 = (centro_x, y_pos)
                pos2 = pos1
        else:
            # Fallback: superior direito
            pos1 = (largura - tamanho - margem, margem)
            pos2 = pos1
        
        return (pos1, pos2)
    
    # =========================================================================
    # GERAÃ‡ÃƒO DE PPTX COM ÃUDIO
    # =========================================================================
    
    def guardar_com_audio(self, caminho_saida: str, incluir_traducao: bool = False,
                          legenda_no_slide: bool = False, legenda_nas_notas: bool = False) -> bool:
        """
        Guarda PPTX com Ã­cones de Ã¡udio nos slides.
        
        Args:
            caminho_saida: Caminho do ficheiro de saÃ­da
            incluir_traducao: Se True, adiciona tambÃ©m Ã­cone para Ã¡udio traduzido
            legenda_no_slide: Se True, adiciona caixa de texto com traduÃ§Ã£o no slide
            legenda_nas_notas: Se True, adiciona traduÃ§Ã£o nas notas do apresentador
        
        Returns:
            True se guardado com sucesso
        """
        if not self.apresentacao or not PPTX_DISPONIVEL:
            return False
        
        try:
            prs = Presentation(self.apresentacao.caminho)
            
            # Criar Ã­cones
            caminho_icone = None
            caminho_icone_trad = None
            if self.config_icone.mostrar:
                caminho_icone = self._criar_icone_audio(64, cor=(41, 128, 185, 255))
                if incluir_traducao:
                    caminho_icone_trad = self._criar_icone_traducao(64, cor=(39, 174, 96, 255))
            
            # Calcular posiÃ§Ãµes dos Ã­cones (considerar se hÃ¡ legenda para nÃ£o sobrepor)
            (pos1, pos2) = self._calcular_posicao_icone(prs, incluir_traducao, tem_legenda=legenda_no_slide)
            tamanho = Cm(self.config_icone.tamanho_cm)
            
            for slide_info in self.apresentacao.slides:
                slide = prs.slides[slide_info.numero - 1]
                
                # Ãcone Ã¡udio principal
                if slide_info.caminho_audio and os.path.exists(slide_info.caminho_audio):
                    if self.config_icone.mostrar and caminho_icone:
                        shape = slide.shapes.add_picture(
                            caminho_icone,
                            pos1[0], pos1[1],
                            width=tamanho, height=tamanho
                        )
                        nome_audio = os.path.basename(slide_info.caminho_audio)
                        shape.click_action.hyperlink.address = nome_audio
                
                # Ãcone Ã¡udio traduzido
                if incluir_traducao and slide_info.caminho_audio_traduzido:
                    if os.path.exists(slide_info.caminho_audio_traduzido):
                        if self.config_icone.mostrar and caminho_icone_trad:
                            shape = slide.shapes.add_picture(
                                caminho_icone_trad,
                                pos2[0], pos2[1],
                                width=tamanho, height=tamanho
                            )
                            nome_audio = os.path.basename(slide_info.caminho_audio_traduzido)
                            shape.click_action.hyperlink.address = nome_audio
                
                # Adicionar caixa de texto com traduÃ§Ã£o no slide
                if legenda_no_slide and slide_info.texto_traduzido.strip():
                    self._adicionar_caixa_texto_traducao(slide, prs, slide_info.texto_traduzido)
                
                # Adicionar traduÃ§Ã£o nas notas
                if legenda_nas_notas and slide_info.texto_traduzido.strip():
                    self._adicionar_traducao_notas(slide, slide_info.texto_traduzido)
            
            # Guardar
            prs.save(caminho_saida)
            
            # v1.9.3: Copiar ficheiros de audio para a pasta do PPTX
            pasta_pptx = os.path.dirname(caminho_saida)
            for slide_info in self.apresentacao.slides:
                # Copiar audio principal
                if slide_info.caminho_audio and os.path.exists(slide_info.caminho_audio):
                    pasta_audio = os.path.dirname(slide_info.caminho_audio)
                    if pasta_audio != pasta_pptx:
                        nome_audio = os.path.basename(slide_info.caminho_audio)
                        destino = os.path.join(pasta_pptx, nome_audio)
                        if not os.path.exists(destino):
                            shutil.copy2(slide_info.caminho_audio, destino)
                
                # Copiar audio traduzido
                if slide_info.caminho_audio_traduzido and os.path.exists(slide_info.caminho_audio_traduzido):
                    pasta_audio = os.path.dirname(slide_info.caminho_audio_traduzido)
                    if pasta_audio != pasta_pptx:
                        nome_audio = os.path.basename(slide_info.caminho_audio_traduzido)
                        destino = os.path.join(pasta_pptx, nome_audio)
                        if not os.path.exists(destino):
                            shutil.copy2(slide_info.caminho_audio_traduzido, destino)
            
            # Limpar icones temporarios
            for ic in [caminho_icone, caminho_icone_trad]:
                if ic and os.path.exists(ic):
                    os.remove(ic)
            
            return True
            
        except Exception as e:
            print(f"Erro ao guardar PPTX: {e}")
            return False
    
    def _adicionar_caixa_texto_traducao(self, slide, prs, texto: str):
        """
        Adiciona caixa de texto com traduÃ§Ã£o na parte inferior do slide.
        Se texto for longo, mostra resumo e adiciona completo nas notas.
        """
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        largura = prs.slide_width
        altura = prs.slide_height
        
        # Tamanho fixo legÃ­vel
        fonte_tamanho = 14
        caixa_altura = Cm(1.8)
        caixa_largura = largura - Cm(2)
        pos_x = Cm(1)
        pos_y = altura - caixa_altura - Cm(0.3)
        
        textbox = slide.shapes.add_textbox(pos_x, pos_y, caixa_largura, caixa_altura)
        tf = textbox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        
        # Limpar texto
        texto_limpo = texto.replace('\n', ' ').strip()
        
        # Limite de caracteres que cabem (~2 linhas)
        max_chars = 180
        
        if len(texto_limpo) > max_chars:
            # Truncar e indicar que hÃ¡ mais nas notas
            texto_mostrar = texto_limpo[:max_chars-20].rsplit(' ', 1)[0] + "... (ver notas)"
            # Adicionar texto completo nas notas
            self._adicionar_traducao_notas(slide, texto)
        else:
            texto_mostrar = texto_limpo
        
        p.text = texto_mostrar
        p.font.size = Pt(fonte_tamanho)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
        
        # Fundo claro
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(248, 248, 248)
    
    def _adicionar_traducao_notas(self, slide, texto: str):
        """Adiciona traduÃ§Ã£o nas notas do apresentador"""
        from pptx.util import Pt
        
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        
        # Adicionar separador e traduÃ§Ã£o
        if notes_tf.text:
            notes_tf.text += "\n\n--- TRADUÃ‡ÃƒO ---\n" + texto
        else:
            notes_tf.text = "--- TRADUÃ‡ÃƒO ---\n" + texto
    
    # =========================================================================
    # EXPORTAÃ‡ÃƒO DE SLIDES PARA IMAGENS
    # =========================================================================
    
    def exportar_slides_imagens(self, pasta_saida: str) -> List[str]:
        """
        Exporta slides como imagens.
        
        Tenta mÃºltiplos mÃ©todos:
        1. LibreOffice + pdftoppm/pdf2image
        2. Imagens placeholder com texto (fallback)
        
        Returns:
            Lista de caminhos das imagens
        """
        if not self.apresentacao:
            return []
        
        os.makedirs(pasta_saida, exist_ok=True)
        
        imagens = []
        
        # MÃ©todo 1: Tentar LibreOffice
        try:
            imagens = self._exportar_via_libreoffice(pasta_saida)
            if imagens:
                return imagens
        except Exception as e:
            print(f"LibreOffice nÃ£o disponÃ­vel: {e}")
        
        # MÃ©todo 2: Criar imagens placeholder com texto
        print("A criar imagens placeholder...")
        imagens = self._criar_imagens_placeholder(pasta_saida)
        
        return imagens
    
    def _exportar_via_libreoffice(self, pasta_saida: str) -> List[str]:
        """Exporta slides via LibreOffice + PDF - v1.9.2: Corrigido para Windows"""
        import platform
        import subprocess
        
        sistema = platform.system()
        pptx_path = os.path.normpath(os.path.abspath(self.apresentacao.caminho))
        pasta_saida = os.path.normpath(os.path.abspath(pasta_saida))
        
        # Determinar caminho do LibreOffice
        soffice = None
        
        if sistema == "Windows":
            # Procurar em Program Files
            program_files = [
                os.environ.get("ProgramFiles", r"C:\Program Files"),
                os.environ.get("ProgramFiles(x86)", r"C:\Program Files (x86)"),
            ]
            for pf in program_files:
                if pf:
                    path = os.path.join(pf, "LibreOffice", "program", "soffice.exe")
                    if os.path.exists(path):
                        soffice = path
                        break
            if not soffice:
                soffice = "soffice.exe"
                
        elif sistema == "Darwin":  # macOS
            soffice_paths = [
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                "/Applications/LibreOffice 7.app/Contents/MacOS/soffice",
                "/Applications/LibreOffice 24.app/Contents/MacOS/soffice",
                os.path.expanduser("~/Applications/LibreOffice.app/Contents/MacOS/soffice"),
            ]
            for path in soffice_paths:
                if os.path.exists(path):
                    soffice = path
                    break
            if not soffice:
                soffice = "soffice"
        else:
            # Linux
            linux_paths = ["/usr/bin/soffice", "/usr/bin/libreoffice", "/usr/local/bin/soffice"]
            for path in linux_paths:
                if os.path.exists(path):
                    soffice = path
                    break
            if not soffice:
                soffice = "soffice"
        
        print(f"[PPTX Handler] LibreOffice: {soffice}")
        print(f"[PPTX Handler] PPTX: {pptx_path}")
        print(f"[PPTX Handler] Pasta saída: {pasta_saida}")
        
        # Converter PPTX para PDF
        cmd_pdf = [
            soffice,
            "--headless",
            "--invisible",
            "--convert-to", "pdf",
            "--outdir", pasta_saida,
            pptx_path
        ]
        
        run_kwargs = {"capture_output": True, "timeout": 180}
        if sistema == "Windows":
            run_kwargs["creationflags"] = subprocess.CREATE_NO_WINDOW
        
        result = subprocess.run(cmd_pdf, **run_kwargs)
        
        if result.returncode != 0:
            erro = result.stderr.decode(errors='ignore') if result.stderr else ""
            print(f"[PPTX Handler] Erro LibreOffice: {erro}")
            raise Exception(f"LibreOffice falhou: {erro}")
        
        # Encontrar PDF gerado
        nome_base = Path(pptx_path).stem
        pdf_path = os.path.join(pasta_saida, f"{nome_base}.pdf")
        
        if not os.path.exists(pdf_path):
            # Tentar encontrar qualquer PDF
            for f in os.listdir(pasta_saida):
                if f.lower().endswith('.pdf'):
                    pdf_path = os.path.join(pasta_saida, f)
                    break
        
        if not os.path.exists(pdf_path):
            raise Exception(f"PDF não gerado em: {pasta_saida}")
        
        print(f"[PPTX Handler] PDF gerado: {pdf_path}")
        
        # Converter PDF para imagens
        imagens = []
        
        # MÉTODO 1: PyMuPDF (recomendado - funciona em Windows sem dependências)
        try:
            import fitz  # PyMuPDF
            print("[PPTX Handler] A usar PyMuPDF")
            
            doc = fitz.open(pdf_path)
            total_pages = len(doc)
            
            for i, pagina in enumerate(doc):
                mat = fitz.Matrix(150/72, 150/72)  # 150 DPI
                pix = pagina.get_pixmap(matrix=mat)
                
                img_path = os.path.join(pasta_saida, f"slide-{i+1:02d}.jpg")
                pix.save(img_path)
                imagens.append(img_path)
                print(f"[PPTX Handler] Slide {i+1}/{total_pages} exportado")
            
            doc.close()
            
            # Remover PDF temporário
            try:
                os.remove(pdf_path)
            except:
                pass
            
            return imagens
            
        except ImportError:
            print("[PPTX Handler] PyMuPDF não disponível")
        except Exception as e:
            print(f"[PPTX Handler] Erro PyMuPDF: {e}")
        
        # MÉTODO 2: pdf2image (precisa Poppler no Windows)
        try:
            from pdf2image import convert_from_path
            print("[PPTX Handler] A usar pdf2image")
            
            pdf2img_kwargs = {"dpi": 150}
            
            if sistema == "Windows":
                # Procurar Poppler
                poppler_paths = [
                    r"C:\Program Files\poppler-24.08.0\Library\bin",
                    r"C:\Program Files\poppler\Library\bin",
                    r"C:\poppler-24.08.0\Library\bin",
                    r"C:\poppler\Library\bin",
                    r"C:\poppler\bin",
                ]
                for pp in poppler_paths:
                    if os.path.exists(os.path.join(pp, "pdftoppm.exe")):
                        pdf2img_kwargs["poppler_path"] = pp
                        print(f"[PPTX Handler] Poppler encontrado: {pp}")
                        break
            
            pages = convert_from_path(pdf_path, **pdf2img_kwargs)
            
            for i, page in enumerate(pages):
                img_path = os.path.join(pasta_saida, f"slide-{i+1:02d}.jpg")
                page.save(img_path, "JPEG", quality=90)
                imagens.append(img_path)
            
            # Remover PDF temporário
            try:
                os.remove(pdf_path)
            except:
                pass
            
            return imagens
            
        except ImportError:
            print("[PPTX Handler] pdf2image não disponível")
        except Exception as e:
            print(f"[PPTX Handler] Erro pdf2image: {e}")
        
        # MÉTODO 3: pdftoppm (Linux/macOS)
        if sistema != "Windows":
            try:
                prefix = os.path.join(pasta_saida, "slide")
                cmd_img = ["pdftoppm", "-jpeg", "-r", "150", pdf_path, prefix]
                subprocess.run(cmd_img, capture_output=True, timeout=120)
                
                for f in sorted(os.listdir(pasta_saida)):
                    if f.startswith("slide") and f.endswith(".jpg"):
                        imagens.append(os.path.join(pasta_saida, f))
                
                if imagens:
                    try:
                        os.remove(pdf_path)
                    except:
                        pass
                    return imagens
            except Exception as e:
                print(f"[PPTX Handler] Erro pdftoppm: {e}")
        
        # Se chegou aqui, nenhum método funcionou
        if not imagens:
            print("\n" + "="*60)
            print("ERRO: Não foi possível converter PDF para imagens.")
            print("SOLUÇÃO: Instale PyMuPDF com o comando:")
            print("    pip install PyMuPDF")
            print("="*60 + "\n")
        
        return imagens
    
    def _criar_imagens_placeholder(self, pasta_saida: str) -> List[str]:
        """
        Cria imagens placeholder com o texto do slide.
        Fallback quando LibreOffice nÃ£o estÃ¡ disponÃ­vel.
        """
        from PIL import Image, ImageDraw, ImageFont
        
        imagens = []
        
        for i, slide in enumerate(self.apresentacao.slides, 1):
            # Criar imagem 1280x720 (16:9)
            img = Image.new('RGB', (1280, 720), color=(45, 45, 60))
            draw = ImageDraw.Draw(img)
            
            # Tentar carregar fonte
            try:
                # Tentar fontes comuns
                font_paths = [
                    "/System/Library/Fonts/Helvetica.ttc",  # macOS
                    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",  # Linux
                    "C:\\Windows\\Fonts\\arial.ttf",  # Windows
                    "arial.ttf"
                ]
                font = None
                for fp in font_paths:
                    try:
                        font = ImageFont.truetype(fp, 32)
                        break
                    except:
                        continue
                if not font:
                    font = ImageFont.load_default()
                
                font_titulo = None
                for fp in font_paths:
                    try:
                        font_titulo = ImageFont.truetype(fp, 48)
                        break
                    except:
                        continue
                if not font_titulo:
                    font_titulo = font
                    
            except:
                font = ImageFont.load_default()
                font_titulo = font
            
            # Desenhar nÃºmero do slide
            draw.text((50, 30), f"Slide {i}", fill=(100, 150, 255), font=font_titulo)
            
            # Obter texto do slide
            texto = slide.texto_narrar if slide.texto_narrar else slide.texto_visivel
            
            if texto:
                # Quebrar texto em linhas
                linhas = self._quebrar_texto_linhas(texto, 60)
                y = 120
                for linha in linhas[:12]:  # MÃ¡ximo 12 linhas
                    draw.text((50, y), linha, fill=(220, 220, 220), font=font)
                    y += 45
                
                if len(linhas) > 12:
                    draw.text((50, y), "...", fill=(150, 150, 150), font=font)
            else:
                draw.text((50, 150), "(Sem texto)", fill=(150, 150, 150), font=font)
            
            # RodapÃ©
            draw.text((50, 670), "PPTX Narrator - Preview", fill=(80, 80, 100), font=font)
            
            # Guardar
            img_path = os.path.join(pasta_saida, f"slide-{i:02d}.jpg")
            img.save(img_path, "JPEG", quality=85)
            imagens.append(img_path)
        
        return imagens
    
    def _quebrar_texto_linhas(self, texto: str, max_chars: int) -> List[str]:
        """Quebra texto em linhas com mÃ¡ximo de caracteres"""
        palavras = texto.replace('\n', ' ').split()
        linhas = []
        linha_atual = ""
        
        for palavra in palavras:
            if len(linha_atual) + len(palavra) + 1 <= max_chars:
                linha_atual = (linha_atual + " " + palavra).strip()
            else:
                if linha_atual:
                    linhas.append(linha_atual)
                linha_atual = palavra
        
        if linha_atual:
            linhas.append(linha_atual)
        
        return linhas
    
    # =========================================================================
    # GERAÃ‡ÃƒO DE FICHEIRO SRT (LEGENDAS)
    # =========================================================================
    
    def gerar_srt(self, caminho_saida: str, usar_traducao: bool = False, 
                  tempo_extra: float = 0.5, tempo_minimo: float = 3.0,
                  max_chars_segmento: int = 120, max_linhas: int = 2) -> bool:
        """
        Gera ficheiro de legendas SRT a partir dos slides.
        v1.7.1: Divide texto em segmentos temporais para nÃ£o encher o ecrÃ£.
        
        Args:
            caminho_saida: Caminho para guardar o ficheiro .srt
            usar_traducao: Se True, usa texto traduzido; senÃ£o, usa texto original
            tempo_extra: Tempo extra apÃ³s cada slide (sincronizar com vÃ­deo)
            tempo_minimo: DuraÃ§Ã£o mÃ­nima para slides sem Ã¡udio
            max_chars_segmento: MÃ¡ximo de caracteres por segmento de legenda
            max_linhas: MÃ¡ximo de linhas por segmento (1-3)
        
        Returns:
            True se gerado com sucesso
        """
        if not self.apresentacao:
            return False
        
        try:
            linhas_srt = []
            tempo_atual = 0.0
            contador = 1
            
            for slide in self.apresentacao.slides:
                # Escolher texto baseado na configuraÃ§Ã£o
                if usar_traducao:
                    texto = slide.texto_traduzido or slide.texto_narrar
                    # Para duraÃ§Ã£o, usar o Ã¡udio que vai ser reproduzido
                    # Se usar traduÃ§Ã£o mas nÃ£o hÃ¡ Ã¡udio traduzido, usar original
                    caminho_audio_trad = slide.caminho_audio_traduzido
                    caminho_audio_orig = slide.caminho_audio
                    
                    # Obter duraÃ§Ãµes reais
                    duracao_trad = self._obter_duracao_audio(caminho_audio_trad) if caminho_audio_trad else 0
                    duracao_orig = self._obter_duracao_audio(caminho_audio_orig) if caminho_audio_orig else 0
                    
                    # Usar a maior duraÃ§Ã£o (para garantir sincronizaÃ§Ã£o)
                    duracao_audio = max(duracao_trad, duracao_orig)
                else:
                    texto = slide.texto_narrar
                    caminho_audio = slide.caminho_audio
                    duracao_audio = self._obter_duracao_audio(caminho_audio) if caminho_audio else 0
                
                # Fallback para duraÃ§Ã£o guardada
                if duracao_audio <= 0:
                    if usar_traducao:
                        duracao_audio = max(slide.duracao_audio_traduzido or 0, slide.duracao_audio or 0)
                    else:
                        duracao_audio = slide.duracao_audio or 0
                
                # Calcular duraÃ§Ã£o total do slide
                if duracao_audio > 0:
                    duracao_slide = duracao_audio + tempo_extra
                else:
                    duracao_slide = tempo_minimo
                
                if not texto or not texto.strip():
                    # Slide sem texto, avanÃ§ar tempo
                    tempo_atual += duracao_slide
                    continue
                
                # Dividir texto em segmentos
                segmentos = self._dividir_texto_segmentos_srt(texto, max_chars_segmento, max_linhas)
                num_segmentos = len(segmentos)
                
                if num_segmentos == 0:
                    tempo_atual += duracao_slide
                    continue
                
                # Calcular duraÃ§Ã£o por segmento
                duracao_por_segmento = duracao_slide / num_segmentos
                
                # Criar entrada SRT para cada segmento
                for seg_idx, segmento in enumerate(segmentos):
                    inicio = tempo_atual + (seg_idx * duracao_por_segmento)
                    fim = inicio + duracao_por_segmento
                    
                    inicio_str = self._formatar_tempo_srt(inicio)
                    fim_str = self._formatar_tempo_srt(fim)
                    
                    # Formatar texto do segmento
                    texto_formatado = self._formatar_texto_srt(segmento, max_chars_linha=60)
                    
                    linhas_srt.append(str(contador))
                    linhas_srt.append(f"{inicio_str} --> {fim_str}")
                    linhas_srt.append(texto_formatado)
                    linhas_srt.append("")
                    
                    contador += 1
                
                tempo_atual += duracao_slide
            
            # Escrever ficheiro
            with open(caminho_saida, 'w', encoding='utf-8') as f:
                f.write('\n'.join(linhas_srt))
            
            return True
            
        except Exception as e:
            print(f"Erro ao gerar SRT: {e}")
            return False
    
    def _dividir_texto_segmentos_srt(self, texto: str, max_chars: int = 120, 
                                      max_linhas: int = 2) -> List[str]:
        """
        Divide texto em segmentos para legendas SRT.
        Cada segmento tem no mÃ¡ximo max_chars caracteres e max_linhas linhas.
        v1.7.1: Nova funÃ§Ã£o para criar legendas paginadas.
        """
        texto = texto.strip().replace('\n', ' ').replace('  ', ' ')
        
        if len(texto) <= max_chars:
            return [texto]
        
        # Dividir por frases primeiro
        import re
        frases = re.split(r'(?<=[.!?])\s+', texto)
        
        segmentos = []
        segmento_atual = ""
        
        for frase in frases:
            frase = frase.strip()
            if not frase:
                continue
            
            # Se a frase cabe no segmento actual
            if len(segmento_atual) + len(frase) + 1 <= max_chars:
                segmento_atual = (segmento_atual + " " + frase).strip()
            else:
                # Guardar segmento actual se nÃ£o estiver vazio
                if segmento_atual:
                    segmentos.append(segmento_atual)
                
                # Se a frase Ã© muito longa, dividir por palavras
                if len(frase) > max_chars:
                    palavras = frase.split()
                    segmento_atual = ""
                    for palavra in palavras:
                        if len(segmento_atual) + len(palavra) + 1 <= max_chars:
                            segmento_atual = (segmento_atual + " " + palavra).strip()
                        else:
                            if segmento_atual:
                                segmentos.append(segmento_atual)
                            segmento_atual = palavra
                else:
                    segmento_atual = frase
        
        # Adicionar Ãºltimo segmento
        if segmento_atual:
            segmentos.append(segmento_atual)
        
        return segmentos if segmentos else [texto[:max_chars]]
    
    def _formatar_texto_srt(self, texto: str, max_chars_linha: int = 60) -> str:
        """
        Formata texto para SRT - quebra em linhas curtas.
        v1.7.1: Reduzido max_chars_linha para legendas mais legÃ­veis.
        """
        texto = texto.strip().replace('\n', ' ').replace('  ', ' ')
        
        if len(texto) <= max_chars_linha:
            return texto
        
        # Quebrar em 2 linhas mÃ¡ximo
        palavras = texto.split()
        linhas = []
        linha_atual = []
        tamanho_atual = 0
        
        for palavra in palavras:
            if tamanho_atual + len(palavra) + 1 <= max_chars_linha:
                linha_atual.append(palavra)
                tamanho_atual += len(palavra) + 1
            else:
                if linha_atual:
                    linhas.append(' '.join(linha_atual))
                linha_atual = [palavra]
                tamanho_atual = len(palavra)
                
                # MÃ¡ximo 2 linhas
                if len(linhas) >= 2:
                    break
        
        if linha_atual and len(linhas) < 2:
            linhas.append(' '.join(linha_atual))
        
        return '\n'.join(linhas[:2])  # MÃ¡ximo 2 linhas
    
    def _formatar_tempo_srt(self, segundos: float) -> str:
        """Formata tempo em segundos para formato SRT (HH:MM:SS,mmm)"""
        horas = int(segundos // 3600)
        minutos = int((segundos % 3600) // 60)
        segs = int(segundos % 60)
        milisseg = int((segundos % 1) * 1000)
        return f"{horas:02d}:{minutos:02d}:{segs:02d},{milisseg:03d}"
    
    def _formatar_texto_legenda(self, texto: str, max_chars: int = 42) -> str:
        """Formata texto para legendas (mÃ¡ximo 2 linhas, 42 chars cada)"""
        texto = texto.strip().replace('\n', ' ').replace('  ', ' ')
        
        if len(texto) <= max_chars:
            return texto
        
        # Dividir em duas linhas
        palavras = texto.split()
        linha1 = []
        linha2 = []
        tamanho_atual = 0
        
        for palavra in palavras:
            if tamanho_atual + len(palavra) + 1 <= max_chars:
                linha1.append(palavra)
                tamanho_atual += len(palavra) + 1
            else:
                linha2.append(palavra)
        
        resultado = ' '.join(linha1)
        if linha2:
            texto_l2 = ' '.join(linha2)
            if len(texto_l2) > max_chars:
                texto_l2 = texto_l2[:max_chars-3] + "..."
            resultado += '\n' + texto_l2
        
        return resultado
    
    def _obter_duracao_audio(self, caminho: str) -> float:
        """
        ObtÃ©m duraÃ§Ã£o de um ficheiro de Ã¡udio em segundos.
        CompatÃ­vel com Python 3.13 (nÃ£o usa audioop/pydub).
        """
        if not caminho or not os.path.exists(caminho):
            return 0.0
        
        # MÃ©todo 1: Usar mutagen (mais fiÃ¡vel)
        try:
            from mutagen.mp3 import MP3
            audio = MP3(caminho)
            return audio.info.length
        except:
            pass
        
        # MÃ©todo 2: Usar ffprobe (ffmpeg)
        try:
            import subprocess
            cmd = [
                "ffprobe", "-v", "quiet", "-show_entries",
                "format=duration", "-of", "default=noprint_wrappers=1:nokey=1",
                caminho
            ]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=10)
            if result.returncode == 0 and result.stdout.strip():
                return float(result.stdout.strip())
        except:
            pass
        
        # MÃ©todo 3: Usar moviepy (se disponÃ­vel)
        try:
            from moviepy import AudioFileClip
            clip = AudioFileClip(caminho)
            duracao = clip.duration
            clip.close()
            return duracao
        except:
            pass
        
        # MÃ©todo 4: Tentar pydub (pode falhar em Python 3.13)
        try:
            from pydub import AudioSegment
            audio = AudioSegment.from_file(caminho)
            return len(audio) / 1000.0
        except:
            pass
        
        return 0.0


# InstÃ¢ncia global
gestor_pptx = GestorPPTX()
